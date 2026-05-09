from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "SxDevOps-AI-Agent-宣传.pptx"
SCREENSHOTS = ROOT / "docs" / "screenshots"


W = Inches(13.333)
H = Inches(7.5)
PAPER_X = Inches(0.38)
PAPER_Y = Inches(0.34)
PAPER_W = Inches(12.58)
PAPER_H = Inches(6.82)
PAD = Inches(0.48)
CONTENT_X = PAPER_X + PAD
CONTENT_Y = PAPER_Y + Inches(0.82)
CONTENT_W = PAPER_W - PAD * 2
CONTENT_H = PAPER_H - Inches(1.42)


COLORS = {
    "bg": RGBColor(23, 27, 31),
    "paper": RGBColor(246, 241, 232),
    "paper_deep": RGBColor(238, 228, 212),
    "ink": RGBColor(32, 36, 42),
    "muted": RGBColor(104, 113, 122),
    "line": RGBColor(216, 207, 191),
    "asset": RGBColor(47, 167, 160),
    "task": RGBColor(214, 139, 54),
    "monitor": RGBColor(79, 124, 207),
    "ai": RGBColor(167, 91, 194),
    "danger": RGBColor(200, 61, 61),
    "ok": RGBColor(63, 155, 98),
    "white": RGBColor(255, 255, 255),
}


FONT = "Microsoft YaHei"
SERIF = "Microsoft YaHei"
MONO = "Consolas"


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]

    # subtle color washes
    wash1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
    wash1.fill.solid()
    wash1.fill.fore_color.rgb = RGBColor(26, 54, 56)
    wash1.fill.transparency = 58
    wash1.line.fill.background()

    paper = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PAPER_X, PAPER_Y, PAPER_W, PAPER_H
    )
    paper.fill.solid()
    paper.fill.fore_color.rgb = COLORS["paper"]
    paper.line.color.rgb = RGBColor(226, 217, 202)
    paper.line.width = Pt(1)
    paper.adjustments[0] = 0.035

    # grid texture
    step = Inches(0.27)
    x = PAPER_X + step
    while x < PAPER_X + PAPER_W - step / 2:
        ln = slide.shapes.add_connector(1, x, PAPER_Y + Inches(0.12), x, PAPER_Y + PAPER_H - Inches(0.12))
        ln.line.color.rgb = RGBColor(232, 224, 210)
        ln.line.transparency = 25
        ln.line.width = Pt(0.35)
        x += step
    y = PAPER_Y + step
    while y < PAPER_Y + PAPER_H - step / 2:
        ln = slide.shapes.add_connector(1, PAPER_X + Inches(0.12), y, PAPER_X + PAPER_W - Inches(0.12), y)
        ln.line.color.rgb = RGBColor(232, 224, 210)
        ln.line.transparency = 30
        ln.line.width = Pt(0.35)
        y += step

    inner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        PAPER_X + Inches(0.12),
        PAPER_Y + Inches(0.12),
        PAPER_W - Inches(0.24),
        PAPER_H - Inches(0.24),
    )
    inner.fill.background()
    inner.line.color.rgb = RGBColor(224, 215, 199)
    inner.line.transparency = 15
    inner.line.width = Pt(0.75)
    inner.adjustments[0] = 0.025


def set_text(box, text, size=18, color="ink", bold=False, font=FONT, align=None, line_spacing=1.05):
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color] if isinstance(color, str) else color
    return box


def add_text(slide, text, x, y, w, h, size=18, color="ink", bold=False, font=FONT, align=None, line_spacing=1.05):
    box = slide.shapes.add_textbox(x, y, w, h)
    return set_text(box, text, size, color, bold, font, align, line_spacing)


def add_topline(slide, slide_no, kicker):
    mark = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PAPER_X + PAD, PAPER_Y + Inches(0.27), Inches(0.36), Inches(0.36)
    )
    mark.fill.solid()
    mark.fill.fore_color.rgb = COLORS["asset"]
    mark.line.fill.background()
    mark.adjustments[0] = 0.16
    add_text(slide, "Sx", PAPER_X + PAD + Inches(0.045), PAPER_Y + Inches(0.335), Inches(0.27), Inches(0.18), 10, "white", True, MONO, PP_ALIGN.CENTER)
    add_text(slide, "SxDevOps AI Agent", PAPER_X + PAD + Inches(0.46), PAPER_Y + Inches(0.31), Inches(3.6), Inches(0.25), 12, "ink", True)
    add_text(slide, kicker.upper(), PAPER_X + PAPER_W - PAD - Inches(3.0), PAPER_Y + Inches(0.32), Inches(3.0), Inches(0.24), 9, "muted", False, MONO, PP_ALIGN.RIGHT)

    y = PAPER_Y + PAPER_H - Inches(0.58)
    ln = slide.shapes.add_connector(1, PAPER_X + PAD, y, PAPER_X + PAPER_W - PAD, y)
    ln.line.color.rgb = COLORS["line"]
    ln.line.width = Pt(1)
    add_text(slide, "可观测性 + 事件墙 + 任务中心 + AIOps", PAPER_X + PAD, y + Inches(0.16), Inches(4.8), Inches(0.22), 9, "muted")
    add_text(slide, f"{slide_no:02d} / 14", PAPER_X + PAPER_W - PAD - Inches(0.8), y + Inches(0.16), Inches(0.8), Inches(0.22), 9, "muted", False, MONO, PP_ALIGN.RIGHT)


def add_title(slide, title, subtitle=None, y=CONTENT_Y, width=CONTENT_W):
    add_text(slide, title, CONTENT_X, y, width, Inches(0.78), 30, "ink", True, SERIF, line_spacing=0.95)
    if subtitle:
        add_text(slide, subtitle, CONTENT_X, y + Inches(0.86), width * 0.78, Inches(0.58), 13, "muted", False, FONT, line_spacing=1.15)


def add_tag(slide, text, x, y, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, Inches(1.22), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(250, 247, 240)
    shape.line.color.rgb = COLORS["line"]
    shape.adjustments[0] = 0.45
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.12), y + Inches(0.12), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = COLORS[color]
    dot.line.fill.background()
    add_text(slide, text, x + Inches(0.32), y + Inches(0.085), Inches(0.82), Inches(0.18), 9, "ink", True)


def add_panel(slide, title, body, x, y, w, h, accent="asset", code=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 252, 246)
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.045
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[accent]
    bar.line.fill.background()
    add_text(slide, title, x + Inches(0.24), y + Inches(0.20), w - Inches(0.42), Inches(0.28), 15, "ink", True)
    add_text(slide, body, x + Inches(0.24), y + Inches(0.62), w - Inches(0.42), h - Inches(0.88), 10.5, "muted", False, FONT, line_spacing=1.18)
    if code:
        add_text(slide, code, x + Inches(0.24), y + h - Inches(0.30), w - Inches(0.42), Inches(0.18), 7.2, accent, False, MONO)


def add_callout(slide, text, x, y, w, h, accent="asset"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 243, 234)
    shape.line.color.rgb = COLORS[accent]
    shape.line.width = Pt(1.2)
    shape.adjustments[0] = 0.035
    add_text(slide, text, x + Inches(0.28), y + Inches(0.20), w - Inches(0.56), h - Inches(0.32), 13, "ink", True, FONT, line_spacing=1.18)


def add_mock(slide, x, y, w, h, title="AI Agent Console"):
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(34, 40, 45)
    frame.line.color.rgb = RGBColor(74, 82, 90)
    frame.adjustments[0] = 0.035
    add_text(slide, title, x + Inches(0.22), y + Inches(0.18), w - Inches(1.1), Inches(0.22), 9.5, "white", True, MONO)
    live = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + w - Inches(0.72), y + Inches(0.15), Inches(0.48), Inches(0.24))
    live.fill.solid()
    live.fill.fore_color.rgb = COLORS["ok"]
    live.line.fill.background()
    live.adjustments[0] = 0.4
    add_text(slide, "LIVE", x + w - Inches(0.64), y + Inches(0.19), Inches(0.32), Inches(0.12), 6.5, "white", True, MONO, PP_ALIGN.CENTER)

    pill_w = (w - Inches(0.58)) / 4
    labels = [("可观测", "monitor"), ("事件", "danger"), ("任务", "task"), ("AI", "ai")]
    for i, (label, color) in enumerate(labels):
        px = x + Inches(0.22) + i * pill_w
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, px, y + Inches(0.58), pill_w - Inches(0.08), Inches(0.34))
        pill.fill.solid()
        pill.fill.fore_color.rgb = COLORS[color] if i == 0 else RGBColor(48, 55, 62)
        pill.line.fill.background()
        pill.adjustments[0] = 0.35
        add_text(slide, label, px, y + Inches(0.665), pill_w - Inches(0.08), Inches(0.12), 7.5, "white", True, FONT, PP_ALIGN.CENTER)

    # KPI blocks
    kpi_w = (w - Inches(0.66)) / 3
    for i, (num, label, color) in enumerate([("12", "活跃告警", "danger"), ("37", "关键事件", "task"), ("4", "待确认动作", "ai")]):
        kx = x + Inches(0.22) + i * kpi_w
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, kx, y + Inches(1.12), kpi_w - Inches(0.08), Inches(0.62))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(44, 51, 57)
        box.line.color.rgb = RGBColor(68, 76, 83)
        box.adjustments[0] = 0.05
        add_text(slide, num, kx + Inches(0.12), y + Inches(1.23), Inches(0.6), Inches(0.18), 17, color, True, MONO)
        add_text(slide, label, kx + Inches(0.77), y + Inches(1.29), kpi_w - Inches(0.92), Inches(0.18), 8.5, "white", False)

    # chart lines
    rows = [("CPU", 0.68, "monitor"), ("MEM", 0.74, "asset"), ("ERR", 0.42, "danger"), ("TASK", 0.54, "task")]
    for idx, (name, pct, color) in enumerate(rows):
        yy = y + Inches(2.05 + idx * 0.36)
        add_text(slide, name, x + Inches(0.28), yy, Inches(0.45), Inches(0.14), 7.2, "white", False, MONO)
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + Inches(0.84), yy + Inches(0.035), w - Inches(1.7), Inches(0.07))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(70, 78, 86)
        bg.line.fill.background()
        fg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + Inches(0.84), yy + Inches(0.035), (w - Inches(1.7)) * pct, Inches(0.07))
        fg.fill.solid()
        fg.fill.fore_color.rgb = COLORS[color]
        fg.line.fill.background()
        add_text(slide, f"{int(pct*100)}%", x + w - Inches(0.64), yy, Inches(0.36), Inches(0.14), 7.2, "white", False, MONO, PP_ALIGN.RIGHT)

    chat = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + Inches(0.22), y + h - Inches(0.94), w - Inches(0.44), Inches(0.62))
    chat.fill.solid()
    chat.fill.fore_color.rgb = RGBColor(27, 32, 37)
    chat.line.color.rgb = RGBColor(63, 72, 80)
    chat.adjustments[0] = 0.06
    add_text(slide, '> "分析 order-center 生产异常，并生成巡检任务草稿"', x + Inches(0.38), y + h - Inches(0.76), w - Inches(0.76), Inches(0.18), 8.2, "white", False, MONO)
    add_text(slide, "已调用：告警 / 日志 / Trace / 事件墙 / 任务中心", x + Inches(0.38), y + h - Inches(0.49), w - Inches(0.76), Inches(0.16), 7.3, "asset", False, MONO)


def add_screenshot(slide, path, x, y, w, h, label=None):
    img_path = Path(path)
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    outer.fill.solid()
    outer.fill.fore_color.rgb = RGBColor(35, 40, 46)
    outer.line.color.rgb = RGBColor(85, 94, 103)
    outer.adjustments[0] = 0.025
    if label:
        add_text(slide, label, x + Inches(0.18), y + Inches(0.13), w - Inches(0.36), Inches(0.18), 8, "white", True, MONO)
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), x + Inches(0.16), y + Inches(0.44), w - Inches(0.32), h - Inches(0.60))


def connect(slide, x1, y1, x2, y2, color="line", width=1.3):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = COLORS[color]
    ln.line.width = Pt(width)


def add_flow_step(slide, idx, title, body, x, y, w, accent):
    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, Inches(0.52), Inches(0.52))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS[accent]
    circle.line.fill.background()
    add_text(slide, f"{idx:02d}", x, y + Inches(0.14), Inches(0.52), Inches(0.16), 9, "white", True, MONO, PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.68), y - Inches(0.02), w - Inches(0.68), Inches(0.26), 14, accent, True)
    add_text(slide, body, x + Inches(0.68), y + Inches(0.35), w - Inches(0.68), Inches(0.56), 10.5, "muted", False, FONT, line_spacing=1.15)


def slide_1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 1, "Unified AI Ops Platform")
    y = CONTENT_Y + Inches(0.22)
    x = CONTENT_X
    add_tag(s, "可观测性", x, y, "monitor")
    add_tag(s, "事件墙", x + Inches(1.36), y, "danger")
    add_tag(s, "任务中心", x + Inches(2.72), y, "task")
    add_tag(s, "AIOps", x + Inches(4.08), y, "ai")
    add_text(s, "SxDevOps\nAI Agent", x, y + Inches(0.72), Inches(5.5), Inches(1.7), 42, "ink", True, SERIF, line_spacing=0.9)
    add_text(s, "智能运维中枢", x, y + Inches(2.28), Inches(4.8), Inches(0.5), 24, "asset", True, SERIF)
    add_text(
        s,
        "把告警、日志、链路、事件、任务和平台资产变成 Agent 可调用的事实工具，"
        "让运维从“到处查系统”升级为“问系统、看证据、确认动作”。",
        x,
        y + Inches(2.92),
        Inches(5.35),
        Inches(0.88),
        14,
        "muted",
        False,
        FONT,
        line_spacing=1.18,
    )
    add_mock(s, CONTENT_X + Inches(6.12), CONTENT_Y + Inches(0.16), Inches(5.55), Inches(4.55))


def slide_2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 2, "Why AI Agent")
    add_title(s, "传统运维的核心问题：信息和动作被拆散了", "排障不是缺少工具，而是缺少能把事实、事件和动作串起来的工作流。")
    x, y = CONTENT_X, CONTENT_Y + Inches(1.62)
    gap = Inches(0.22); w = (CONTENT_W - gap * 3) / 4; h = Inches(2.15)
    panels = [
        ("告警散", "告警中心看到红点，但缺少日志、Trace、最近变更和资产上下文。", "danger"),
        ("事件散", "发布、执行、失败、审批、关键写操作散落在不同页面，难以复盘。", "task"),
        ("任务散", "巡检、命令、Playbook 与主机权限脱节，动作入口不统一。", "asset"),
        ("经验散", "排障路径依赖个人记忆，结论难沉淀，复用成本高。", "ai"),
    ]
    for i, item in enumerate(panels):
        add_panel(s, item[0], item[1], x + i * (w + gap), y, w, h, item[2])
    add_callout(s, "SxDevOps AI Agent 的定位不是再做一个聊天框，而是把平台已有能力包装成可控、可审计、可确认的智能工作流。", x, y + h + Inches(0.45), CONTENT_W, Inches(0.8), "ai")


def slide_3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 3, "Product Map")
    add_title(s, "SxDevOps AI Agent = 四个核心能力的统一入口", "Agent 不是替代平台页面，而是把平台能力转成自然语言可调度的运维接口。")
    y = CONTENT_Y + Inches(1.55)
    steps = [
        ("可观测性", "聚合告警、日志、链路追踪、Grafana 大屏，形成一条连续排障路径。", "monitor"),
        ("事件墙", "聚焦最终执行结果与关键写操作，按业务线、环境、应用快速收敛故障范围。", "danger"),
        ("任务中心", "查询最近任务、生成主机巡检/服务检查/批量命令任务草稿，确认后执行。", "task"),
        ("AIOps", "LLM 负责理解与规划，MCP 提供事实，Skill 约束输出，后端负责安全与审计。", "ai"),
    ]
    x = CONTENT_X
    for i, (title, body, color) in enumerate(steps, 1):
        add_flow_step(s, i, title, body, x, y + Inches((i - 1) * 0.96), CONTENT_W, color)
        if i < len(steps):
            connect(s, x + Inches(0.26), y + Inches((i - 1) * 0.96 + 0.52), x + Inches(0.26), y + Inches(i * 0.96), color, 1.2)
    add_callout(s, "一句话：可观测性负责取证，事件墙负责复盘，任务中心负责行动，AIOps 负责把它们组织成一次可解释的运维处理链路。", CONTENT_X + Inches(6.35), CONTENT_Y + Inches(1.7), Inches(5.35), Inches(2.2), "asset")


def slide_4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 4, "Observability")
    add_title(s, "可观测性平台：从告警入口走到证据闭环", "日志、告警、Trace、Grafana 不再是四个孤立入口，而是 Agent 取证的统一事实层。")
    add_panel(s, "统一入口", "平台总览聚合日志、告警、Trace、Grafana 摘要、入口和最近活动。", CONTENT_X, CONTENT_Y + Inches(1.54), Inches(2.95), Inches(1.35), "monitor")
    add_panel(s, "联动排障", "告警可跳转日志、链路追踪和 Grafana，看同一条件下的上下文证据。", CONTENT_X, CONTENT_Y + Inches(3.05), Inches(2.95), Inches(1.35), "asset")
    add_panel(s, "Agent 工具化", "query_observability / query_alerts / query_logs / query_traces 变成模型可调用工具。", CONTENT_X, CONTENT_Y + Inches(4.56), Inches(2.95), Inches(1.05), "ai")
    add_screenshot(s, SCREENSHOTS / "logs-or-sql.png", CONTENT_X + Inches(3.35), CONTENT_Y + Inches(1.5), Inches(8.3), Inches(4.18), "Observability Console")


def slide_5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 5, "Event Wall")
    add_title(s, "事件墙：把关键结果沉淀成可追溯现场", "事件墙不是流水账，而是围绕最终执行结果、关键写操作和失败定位的复盘面板。")
    add_screenshot(s, SCREENSHOTS / "event-wall.png", CONTENT_X, CONTENT_Y + Inches(1.48), Inches(6.86), Inches(4.12), "Event Wall")
    x = CONTENT_X + Inches(7.15)
    add_panel(s, "失败定位", "按业务线、环境、应用过滤，快速看失败事件集中在哪个范围。", x, CONTENT_Y + Inches(1.48), Inches(4.52), Inches(1.05), "danger")
    add_panel(s, "操作审计", "只保留最终执行结果和关键写操作，默认过滤未执行的驳回审批流。", x, CONTENT_Y + Inches(2.78), Inches(4.52), Inches(1.05), "task")
    add_panel(s, "Agent 复盘", "query_event_wall 能回答“最近生产失败了什么”“谁触发了关键变更”。", x, CONTENT_Y + Inches(4.08), Inches(4.52), Inches(1.05), "ai")


def slide_6(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 6, "Task Center")
    add_title(s, "任务中心：让 Agent 从“会分析”走向“可行动”", "动作型能力必须先生成任务草稿，再由用户确认，最终进入任务中心执行和审计。")
    y = CONTENT_Y + Inches(1.55)
    labels = [
        ("查询任务", "最近巡检、部署、命令执行、失败任务与耗时统计", "monitor"),
        ("生成草稿", "自动补齐目标主机、命令、超时、风险提示和执行策略", "ai"),
        ("二次确认", "高风险动作由用户确认，避免误创建、误执行", "task"),
        ("落入中心", "确认后写入任务中心，执行结果进入事件墙和审计", "asset"),
    ]
    x = CONTENT_X
    w = (CONTENT_W - Inches(0.72)) / 4
    for i, (t, b, c) in enumerate(labels):
        add_panel(s, t, b, x + i * (w + Inches(0.24)), y, w, Inches(1.95), c)
        if i < 3:
            connect(s, x + i * (w + Inches(0.24)) + w, y + Inches(0.98), x + (i + 1) * (w + Inches(0.24)), y + Inches(0.98), c, 1.4)
    add_callout(s, "核心边界：咨询类问题直接返回事实；生成类问题先出草稿；执行类动作必须确认后才落任务。", CONTENT_X, y + Inches(2.42), CONTENT_W, Inches(0.92), "task")


def slide_7(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 7, "AIOps Agent")
    add_title(s, "AIOps：不是关键词查表，而是工具调用型智能体", "LLM 负责理解和规划，平台负责工具、权限、安全边界、审计和兜底。")
    x, y = CONTENT_X, CONTENT_Y + Inches(1.52)
    gap = Inches(0.26); w = (CONTENT_W - gap * 2) / 3
    add_panel(s, "MCP 工具层", "CMDB、可观测性、事件墙、任务中心、容器、中间件等平台能力统一暴露为工具 schema。", x, y, w, Inches(2.15), "monitor", "query_observability · query_event_wall")
    add_panel(s, "Skill 整形层", "约束回答结构，让告警分析、工单汇总、任务生成有稳定输出模板。", x + w + gap, y, w, Inches(2.15), "ai", "结论 / 依据 / 建议操作")
    add_panel(s, "安全执行层", "工具白名单、参数清洗、RBAC、超时、异常兜底、动作确认和审计记录。", x + 2 * (w + gap), y, w, Inches(2.15), "task", "RBAC · Pending Action · Audit")
    add_callout(s, "当前实现本质：用 LLM 做理解和规划，用 MCP 提供事实能力，用 Skill 约束回答结构，用平台后端负责安全、审计、执行与兜底。", x, y + Inches(2.62), CONTENT_W, Inches(1.0), "ai")


def slide_8(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 8, "Two-stage Answer")
    add_title(s, "一次问答链路：先取证，再整形，最后可审计", "让最终回答稳定、可读、可追溯，不因为模型自由发挥丢失工具事实。")
    steps = [
        ("用户提问", "例如：分析生产 order-center 最近异常", "asset"),
        ("LLM 规划", "根据工具 schema 选择告警、日志、Trace、事件墙、任务工具", "ai"),
        ("平台执行", "后端按白名单执行工具，做权限、参数清洗、超时控制", "monitor"),
        ("事实汇总", "结构化工具结果形成证据集和代码兜底草稿", "task"),
        ("Skill 整形", "按模板输出结论、依据、建议操作和可继续查看项", "ai"),
    ]
    y = CONTENT_Y + Inches(1.46)
    for i, (t, b, c) in enumerate(steps, 1):
        yy = y + Inches((i - 1) * 0.75)
        add_flow_step(s, i, t, b, CONTENT_X, yy, CONTENT_W, c)
    add_callout(s, "模型不直接访问数据库、不直接执行危险动作；它只能提出工具调用计划，实际执行始终由平台后端控制。", CONTENT_X + Inches(7.1), CONTENT_Y + Inches(1.64), Inches(4.55), Inches(1.9), "danger")


def slide_9(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 9, "Scenario")
    add_title(s, "典型场景：生产 order-center 异常分析", "把告警、日志、链路、事件、任务串成一次完整的处理过程。")
    x, y = CONTENT_X, CONTENT_Y + Inches(1.48)
    add_panel(s, "Before：人工排障", "1. 打开告警平台看标签\n2. 去日志中心查错误\n3. 到链路追踪找慢调用\n4. 翻最近发布和事件\n5. 手工写结论和处理建议", x, y, Inches(5.48), Inches(3.5), "danger")
    add_panel(s, "After：AI Agent 处理", "1. 自然语言发起分析\n2. 自动调用告警 / 日志 / Trace / 事件墙\n3. 输出结论、依据、风险和建议操作\n4. 需要巡检时生成任务草稿\n5. 用户确认后进入任务中心", x + Inches(6.1), y, Inches(5.48), Inches(3.5), "ai")
    add_callout(s, "效果：排障链路从“找资料”压缩为“看证据、做判断、确认动作”。", x, y + Inches(3.92), CONTENT_W, Inches(0.72), "asset")


def slide_10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 10, "Safety")
    add_title(s, "安全护栏：让智能体可用，也可控", "企业运维场景里，Agent 的边界比能力更重要。")
    x, y = CONTENT_X, CONTENT_Y + Inches(1.45)
    items = [
        ("RBAC 权限", "后端接口、前端路由、菜单和 WebSocket 场景统一做权限收敛。", "asset"),
        ("工具白名单", "模型只能调用平台注册过的工具，真正执行由后端完成。", "monitor"),
        ("动作确认", "任务生成先进入待确认动作，用户确认后才写入任务中心。", "task"),
        ("审计回放", "会话、工具调用、待确认动作和关键事件都可追踪。", "ai"),
    ]
    gap = Inches(0.24); w = (CONTENT_W - gap) / 2
    for i, (t, b, c) in enumerate(items):
        add_panel(s, t, b, x + (i % 2) * (w + gap), y + (i // 2) * Inches(1.62), w, Inches(1.32), c)
    add_callout(s, "原则：回答基于事实，不凭空编造；动作经过确认，不越权执行；过程留下记录，不做黑盒运维。", x, y + Inches(3.62), CONTENT_W, Inches(0.82), "danger")


def slide_11(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 11, "Architecture")
    add_title(s, "平台架构：数据进来，Agent 编排，动作闭环", "基于 Django / DRF / Channels 与 Vue 3，把平台工程能力产品化。")
    y = CONTENT_Y + Inches(1.55)
    cols = [
        ("数据与事实层", "CMDB 资产\n告警 / 日志 / Trace\n事件墙\n任务中心\n容器 / 中间件", "monitor"),
        ("智能体编排层", "Tool-calling\n内置 MCP\n外部 MCP\nSkill 模板\n兜底草稿", "ai"),
        ("安全控制层", "RBAC\n参数清洗\n工具白名单\n动作确认\n审计记录", "task"),
        ("体验入口层", "右下角浮窗\nAIOps 配置\n会话历史\n思考过程\n动作卡片", "asset"),
    ]
    w = (CONTENT_W - Inches(0.72)) / 4
    for i, (t, b, c) in enumerate(cols):
        add_panel(s, t, b, CONTENT_X + i * (w + Inches(0.24)), y, w, Inches(3.65), c)
        if i < 3:
            connect(s, CONTENT_X + i * (w + Inches(0.24)) + w, y + Inches(1.82), CONTENT_X + (i + 1) * (w + Inches(0.24)), y + Inches(1.82), c, 1.3)


def slide_12(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 12, "Value")
    add_title(s, "SxDevOps AI Agent 带来的产品价值", "从工具集合升级为围绕事实、事件和动作的智能运维系统。")
    y = CONTENT_Y + Inches(1.7)
    metrics = [
        ("可见", "告警、日志、Trace、Grafana、资产状态一眼可查。", "monitor"),
        ("可追", "事件墙记录最终结果、关键写操作和失败定位线索。", "danger"),
        ("可执行", "任务中心承接巡检、命令、Playbook 与批量操作。", "task"),
        ("可审计", "会话、工具调用、待确认动作、执行结果全链路留痕。", "ai"),
    ]
    w = (CONTENT_W - Inches(0.72)) / 4
    for i, (word, body, color) in enumerate(metrics):
        x = CONTENT_X + i * (w + Inches(0.24))
        shape = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, Inches(2.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 252, 246)
        shape.line.color.rgb = COLORS["line"]
        shape.adjustments[0] = 0.05
        add_text(s, word, x, y + Inches(0.36), w, Inches(0.55), 28, color, True, SERIF, PP_ALIGN.CENTER)
        add_text(s, body, x + Inches(0.24), y + Inches(1.22), w - Inches(0.48), Inches(0.72), 10.5, "muted", False, FONT, PP_ALIGN.CENTER, line_spacing=1.12)
    add_callout(s, "最终效果：少切系统，少拼证据，少靠口头经验，多用事实驱动、多用确认动作闭环。", CONTENT_X, y + Inches(2.82), CONTENT_W, Inches(0.82), "asset")


def slide_13(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 13, "Roadmap")
    add_title(s, "下一步：从可观测智能体走向运维执行智能体", "先把事实链路做可信，再把高频、低风险动作纳入自动化闭环。")
    x, y = CONTENT_X, CONTENT_Y + Inches(1.55)
    add_panel(s, "模板族", "告警处置、工单汇总、成本分析、K8s 异常、任务生成等问题类型拥有独立 Skill 模板。", x, y, Inches(3.62), Inches(2.0), "ai")
    add_panel(s, "MCP 扩展", "继续扩展平台内置 MCP，同时增强外部 MCP 健康检查、工具发现、鉴权与超时诊断。", x + Inches(4.0), y, Inches(3.62), Inches(2.0), "monitor")
    add_panel(s, "处置编排", "在只读诊断后，接入审批、命令模板、Runbook 和任务编排，形成低风险自动化闭环。", x + Inches(8.0), y, Inches(3.62), Inches(2.0), "task")
    add_callout(s, "演进路径：先可信取证，再标准输出，最后把可控动作接入任务中心。", x, y + Inches(2.5), CONTENT_W, Inches(0.88), "asset")


def slide_14(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s); add_topline(s, 14, "Close")
    add_text(s, "SxDevOps AI Agent", CONTENT_X, CONTENT_Y + Inches(1.1), CONTENT_W, Inches(0.72), 36, "ink", True, SERIF, PP_ALIGN.CENTER)
    add_text(s, "让运维从“查系统”升级为\n“问系统、看证据、确认动作”", CONTENT_X + Inches(1.0), CONTENT_Y + Inches(2.05), CONTENT_W - Inches(2.0), Inches(1.25), 30, "asset", True, SERIF, PP_ALIGN.CENTER, line_spacing=0.98)
    tag_y = CONTENT_Y + Inches(3.78)
    total_w = Inches(5.3)
    start = CONTENT_X + (CONTENT_W - total_w) / 2
    add_tag(s, "可观测性", start, tag_y, "monitor")
    add_tag(s, "事件墙", start + Inches(1.36), tag_y, "danger")
    add_tag(s, "任务中心", start + Inches(2.72), tag_y, "task")
    add_tag(s, "AIOps", start + Inches(4.08), tag_y, "ai")
    add_text(s, "一个面向真实运维现场的智能运维 Agent 项目", CONTENT_X, CONTENT_Y + Inches(4.55), CONTENT_W, Inches(0.3), 14, "muted", False, FONT, PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    for fn in [
        slide_1,
        slide_2,
        slide_3,
        slide_4,
        slide_5,
        slide_6,
        slide_7,
        slide_8,
        slide_9,
        slide_10,
        slide_11,
        slide_12,
        slide_13,
        slide_14,
    ]:
        fn(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
